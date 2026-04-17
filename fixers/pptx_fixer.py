# fixers/pptx_fixer.py

import io
import logging
import warnings
from lxml import etree
from PIL import Image

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

import torch
from torchvision.models import mobilenet_v3_small, MobileNet_V3_Small_Weights

import wcag_contrast_ratio as wcag

logging.getLogger("transformers").setLevel(logging.ERROR)
warnings.filterwarnings("ignore", category=UserWarning)

# ── XML namespaces ────────────────────────────────────────────────────────────
# DrawingML (colours, fills, fonts)
_DML = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
# PresentationML (slides, shapes, backgrounds)
_PML = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# Maps <a:schemeClr val="…"> semantic aliases → clrScheme child-element names
_SCHEME_ALIAS = {
    "bg1":   "lt1",
    "bg2":   "lt2",
    "tx1":   "dk1",
    "tx2":   "dk2",
    "phClr": "dk1",   # placeholder colour – fall back to dark-1
}

# ============================================================
# MOBILENET
# ============================================================

_MN_MODEL      = None
_MN_TRANSFORM  = None
_MN_CATEGORIES = None


def _load_mobilenet():
    global _MN_MODEL, _MN_TRANSFORM, _MN_CATEGORIES
    if _MN_MODEL is None:
        weights        = MobileNet_V3_Small_Weights.DEFAULT
        _MN_MODEL      = mobilenet_v3_small(weights=weights)
        _MN_MODEL.eval()
        _MN_TRANSFORM  = weights.transforms()
        _MN_CATEGORIES = weights.meta["categories"]


def classify_image(image_bytes):
    _load_mobilenet()
    img   = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    batch = _MN_TRANSFORM(img).unsqueeze(0)
    with torch.no_grad():
        preds = _MN_MODEL(batch)
        probs = preds.softmax(dim=1)[0]
        score, idx = probs.max(dim=0)
    return _MN_CATEGORIES[int(idx)], float(score)


def is_decorative_label(label):
    decorative_keywords = [
        "pattern", "background", "wallpaper", "border", "frame",
        "logo", "emblem", "symbol", "texture", "screen", "monitor",
        "banner", "sign", "flag",
    ]
    return any(k in label.lower() for k in decorative_keywords)


# ============================================================
# THEME COLOUR EXTRACTION
# ============================================================

def extract_theme_colors(prs):
    """
    Walk the first slide-master's relationships to find the theme part,
    parse its <a:clrScheme>, and return:
        { 'dk1': 'RRGGBB', 'lt1': 'RRGGBB', 'acc1': 'RRGGBB', … }
    Returns {} on any failure so callers fall back gracefully.
    """
    try:
        slide_master = prs.slide_masters[0]
        theme_part   = None
        for rel in slide_master.part.rels.values():
            if "theme" in rel.reltype.lower():
                theme_part = rel.target_part
                break
        if theme_part is None:
            return {}

        tree       = etree.fromstring(theme_part.blob)
        clr_scheme = tree.find(f".//{_DML}clrScheme")
        if clr_scheme is None:
            return {}

        colors = {}
        for child in clr_scheme:
            tag  = child.tag.split("}")[-1]          # e.g. "dk1", "lt2", "acc1"
            srgb = child.find(f"{_DML}srgbClr")
            if srgb is not None:
                colors[tag] = srgb.get("val", "000000")
                continue
            sys_clr = child.find(f"{_DML}sysClr")
            if sys_clr is not None:
                colors[tag] = sys_clr.get("lastClr", "000000")
        return colors
    except Exception:
        return {}


# ============================================================
# LOW-LEVEL XML COLOUR RESOLUTION
# ============================================================

def _hex_to_rgb(hex_str):
    s = (hex_str or "000000").lstrip("#").zfill(6)
    try:
        return tuple(int(s[i: i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return (0, 0, 0)


def _apply_lum_modifiers(rgb, lum_mod_val, lum_off_val):
    """Apply <a:lumMod> / <a:lumOff> (integer 1/100 000 fractions)."""
    mod = lum_mod_val / 100_000.0
    off = lum_off_val / 100_000.0

    def adjust(c):
        return max(0, min(255, round(c * mod + off * 255)))

    return tuple(adjust(c) for c in rgb)


def resolve_color_element(elem, theme_colors):
    """
    Given an XML colour child element (direct child of <a:solidFill>)
    return (R, G, B) or None.

    Handles:
      <a:srgbClr val="RRGGBB"/>
      <a:sysClr  lastClr="RRGGBB"/>
      <a:schemeClr val="dk1">
          <a:lumMod val="75000"/>
          <a:lumOff val="25000"/>
      </a:schemeClr>
    """
    if elem is None:
        return None

    local = elem.tag.split("}")[-1]

    if local == "srgbClr":
        return _hex_to_rgb(elem.get("val", "000000"))

    if local == "sysClr":
        return _hex_to_rgb(elem.get("lastClr", "000000"))

    if local == "schemeClr":
        scheme_key   = elem.get("val", "")
        resolved_key = _SCHEME_ALIAS.get(scheme_key, scheme_key)
        hex_val      = theme_colors.get(resolved_key, "000000")
        rgb          = _hex_to_rgb(hex_val)

        lm      = elem.find(f"{_DML}lumMod")
        lo      = elem.find(f"{_DML}lumOff")
        lum_mod = int(lm.get("val", "100000")) if lm is not None else 100_000
        lum_off = int(lo.get("val", "0"))       if lo is not None else 0

        if lum_mod != 100_000 or lum_off != 0:
            rgb = _apply_lum_modifiers(rgb, lum_mod, lum_off)

        return rgb

    return None


def _resolve_solid_fill(solid_elem, theme_colors):
    """Resolve every colour child of a <a:solidFill>; return first success."""
    if solid_elem is None:
        return None
    for child in solid_elem:
        rgb = resolve_color_element(child, theme_colors)
        if rgb is not None:
            return rgb
    return None


# ============================================================
# TARGETED FILL READERS
# (never touch text-body / run children)
# ============================================================

def _spPr_fill(shape_elem, theme_colors):
    """
    Read the fill from a shape's <p:spPr> or <a:spPr> ONLY.

    This is a *direct-child* search — it will never descend into
    <p:txBody> or any <a:rPr> that lives inside text runs, so it
    cannot accidentally return a text colour as a background colour.

    Returns:
        (R, G, B)   – explicit solid fill
        "noFill"    – explicit transparent fill
        None        – no fill information found
    """
    # Shapes use the PML namespace; group/picture shapes may use DML
    sp_pr = shape_elem.find(f"{_PML}spPr")
    if sp_pr is None:
        sp_pr = shape_elem.find(f"{_DML}spPr")
    if sp_pr is None:
        return None

    # Explicit noFill → treat as transparent (use slide background)
    if sp_pr.find(f"{_DML}noFill") is not None:
        return "noFill"

    # Direct-child solidFill only (no .// wildcard)
    solid = sp_pr.find(f"{_DML}solidFill")
    return _resolve_solid_fill(solid, theme_colors)


def _slide_bg_fill(slide, theme_colors):
    """
    Read the solid fill from the slide's <p:bg><p:bgPr>.

    Only inspects the background property element directly;
    never descends into shape or text children.

    Returns (R, G, B) or None.
    """
    try:
        # slide.background._element is <p:bg>
        bg_elem = slide.background._element
        bg_pr   = bg_elem.find(f"{_PML}bgPr")
        if bg_pr is not None:
            if bg_pr.find(f"{_DML}noFill") is not None:
                return None
            solid = bg_pr.find(f"{_DML}solidFill")
            rgb   = _resolve_solid_fill(solid, theme_colors)
            if rgb is not None:
                return rgb
    except Exception:
        pass
    return None


# ============================================================
# SHAPE / RUN COLOUR GETTERS
# ============================================================

def get_shape_bg_color(shape, slide, theme_colors):
    """
    Effective background colour behind *shape*.  Resolution order:

      1. Shape's own <p:spPr> solidFill  (direct child — never text runs)
      2. If shape is transparent (noFill) or has no fill → slide <p:bg>
      3. Default white

    The previous implementation used shape._element.find('.//{_DML}solidFill')
    which searched ALL descendants, including <a:rPr> nodes inside text runs.
    That caused fg == bg (ratio 1.00) false positives, and made fixes appear
    to introduce new violations because the newly-written run solidFill nodes
    were then picked up as the background on re-check.
    """
    result = _spPr_fill(shape._element, theme_colors)

    # Explicit solid fill found on the shape itself
    if result is not None and result != "noFill":
        return result

    # Shape is transparent or inherits — fall through to slide background
    slide_bg = _slide_bg_fill(slide, theme_colors)
    if slide_bg is not None:
        return slide_bg

    return (255, 255, 255)     # safe default: assume white canvas


def get_run_fg_color(run, theme_colors):
    """
    Foreground (font) colour of *run*.  Resolution order:
      1. <a:solidFill> that is a direct child of the run's <a:rPr> (XML)
      2. python-pptx high-level API (handles plain RGB already resolved)
      3. Default black
    """
    try:
        rPr = run._r.find(f"{_DML}rPr")
        if rPr is not None:
            # Direct child only — rPr has no sub-elements that would confuse us
            solid = rPr.find(f"{_DML}solidFill")
            rgb   = _resolve_solid_fill(solid, theme_colors)
            if rgb is not None:
                return rgb
    except Exception:
        pass

    try:
        cf = run.font.color
        if cf.type is not None and cf.rgb is not None:
            rgb = cf.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass

    return (0, 0, 0)


# ============================================================
# CONTRAST HELPERS
# ============================================================

def contrast_ratio(c1, c2):
    f1 = tuple(v / 255.0 for v in c1)
    f2 = tuple(v / 255.0 for v in c2)
    return wcag.rgb(f1, f2)


def is_large_text(run):
    try:
        if run.font.size:
            pt = run.font.size.pt
            return pt >= 18 or (pt >= 14 and run.font.bold)
    except Exception:
        pass
    return False


def fix_run_contrast_xml(run, bg):
    """
    Fix the run's font colour by directly rewriting its <a:rPr> XML.

    Using run.font.color.rgb = … via the python-pptx API only appends a new
    srgbClr node but leaves the original <a:schemeClr> in place; many
    renderers continue to honour the theme colour.  Instead we:
      1. Locate or create <a:rPr> as the first child of the run element.
      2. Remove every existing <a:solidFill> (may contain <a:schemeClr>,
         <a:srgbClr>, or <a:sysClr>).
      3. Insert a fresh <a:solidFill><a:srgbClr val="RRGGBB"/></a:solidFill>
         so the run has one authoritative, explicit colour.

    Picks whichever of black / white gives better contrast against *bg*.
    """
    black = (0, 0, 0)
    white = (255, 255, 255)
    target  = black if contrast_ratio(black, bg) >= contrast_ratio(white, bg) else white
    hex_val = "{:02X}{:02X}{:02X}".format(*target)

    r_elem = run._r

    # Locate or create <a:rPr> — must be first child of <a:r>
    rPr = r_elem.find(f"{_DML}rPr")
    if rPr is None:
        rPr = etree.Element(f"{_DML}rPr")
        r_elem.insert(0, rPr)

    # Remove every existing solidFill so nothing is inherited or duplicated
    for sf in rPr.findall(f"{_DML}solidFill"):
        rPr.remove(sf)

    # Insert fresh solidFill at position 0 (before <a:latin>, <a:ea>, etc.)
    solid = etree.Element(f"{_DML}solidFill")
    srgb  = etree.SubElement(solid, f"{_DML}srgbClr")
    srgb.set("val", hex_val)
    rPr.insert(0, solid)


# ============================================================
# IMAGE HELPERS
# ============================================================

def extract_image_bytes(shape):
    try:
        return shape.image.blob
    except Exception:
        return None


def fix_image_alt(shape, issues, slide_idx):
    image_bytes = extract_image_bytes(shape)
    if not image_bytes:
        return 0, 0

    label, score = classify_image(image_bytes)

    cNvPr = None
    for e in shape._element.iter():
        if e.tag.endswith("cNvPr"):
            cNvPr = e
            break
    if cNvPr is None:
        return 0, 0

    descr = (cNvPr.get("descr") or "").strip()
    title = (cNvPr.get("title") or "").strip()

    if is_decorative_label(label):
        cNvPr.set("title", "Decorative")
        cNvPr.set("descr", "")
        issues.append(("Decorative image normalised", f"Slide {slide_idx + 1}"))
        return 0, 1

    if not descr and not title:
        cNvPr.set("descr", f"Image of {label}")
        issues.append(("Alt text added", f"Slide {slide_idx + 1}"))
        return 1, 0

    return 0, 0


# ============================================================
# TABLE HEADER HELPERS
# ============================================================

def check_and_fix_table_headers(shape, issues, slide_idx, apply_fix=False):
    """
    Verifies that a table shape has its 'first row as header' flag set
    (<a:tblPr firstRow="1"/>), exposed by python-pptx as table.first_row.
    When apply_fix is True also bolds first-row cells for visual clarity.
    Returns the number of violations found (0 or 1 per table).
    """
    violations = 0
    try:
        table = shape.table
        if not table.first_row:
            violations = 1
            issues.append(("Table missing header row flag", f"Slide {slide_idx + 1}"))

            if apply_fix:
                table.first_row = True
                if len(table.rows) > 0:
                    for cell in table.rows[0].cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.bold = True
                issues.append((
                    "Table header row flag set + first row bolded",
                    f"Slide {slide_idx + 1}",
                ))
    except Exception:
        pass
    return violations


# ============================================================
# MAIN SLIDE PROCESSOR
# ============================================================

def process_slides(prs, apply_fix=False):
    """
    Iterate every slide and shape, checking / optionally fixing:
      - Image alt text
      - Table header flags
      - Text colour contrast

    Returns a 6-tuple:
        missing_alt         (int)
        decorative          (int)
        contrast_violations (int)
        total_text_runs     (int)  – denominator for proportional contrast score
        table_violations    (int)
        issues              (list of (description, location) tuples)
    """
    missing_alt         = 0
    decorative          = 0
    contrast_violations = 0
    total_text_runs     = 0
    table_violations    = 0
    issues              = []

    theme_colors = extract_theme_colors(prs)

    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:

            # -------- PICTURES --------
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                m, d = fix_image_alt(shape, issues, s_idx)
                missing_alt += m
                decorative  += d

            # -------- TABLES --------
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_violations += check_and_fix_table_headers(
                    shape, issues, s_idx, apply_fix=apply_fix
                )

            # -------- TEXT --------
            if getattr(shape, "has_text_frame", False):
                bg = get_shape_bg_color(shape, slide, theme_colors)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue

                        total_text_runs += 1
                        fg        = get_run_fg_color(run, theme_colors)
                        threshold = 3.0 if is_large_text(run) else 4.5
                        ratio     = contrast_ratio(fg, bg)

                        if ratio < threshold:
                            contrast_violations += 1
                            issues.append((
                                f"Low contrast text (ratio {ratio:.2f}, "
                                f"threshold {threshold})",
                                f"Slide {s_idx + 1}",
                            ))
                            if apply_fix:
                                fix_run_contrast_xml(run, bg)

    return (
        missing_alt,
        decorative,
        contrast_violations,
        total_text_runs,
        table_violations,
        issues,
    )


# ============================================================
# LANGUAGE FIX
# ============================================================

def fix_language(prs, issues):
    if not prs.core_properties.language:
        prs.core_properties.language = "en-US"
        issues.append(("Language set to en-US", "Presentation"))
        return 1
    return 0


# ============================================================
# SCORING
# ============================================================

class AllyScores:

    @staticmethod
    def score_from_count(count):
        """Step-function penalty for alt-text, decorative, tables, lists, links."""
        if count == 0:
            return 100
        if count <= 5:
            return 76
        if count <= 21:
            return 53
        return 5

    @staticmethod
    def contrast_score(violations, total):
        """
        Proportional Ally-style score for colour contrast.

        Ally measures the percentage of content that passes and maps it to a
        0-100 score.  A floor of 25 prevents a single bad run from collapsing
        an otherwise healthy deck.

        Examples (violations / total → score):
          0  / 100 → 100   (perfect)
          5  / 100 →  95
          20 / 100 →  80
          50 / 100 →  50
         100 / 100 →  25   (floor)
        """
        if total == 0 or violations == 0:
            return 100
        passing_ratio = (total - violations) / total
        return max(25, round(passing_ratio * 100))

    @staticmethod
    def compute(
        missing_alt,
        decorative,
        language_missing,
        headings_missing,
        tables_missing,
        contrast,
        total_text_runs,
        lists,
        links,
    ):
        scores = {
            "alternative_text": AllyScores.score_from_count(missing_alt),
            "decorative":        AllyScores.score_from_count(decorative),
            "language":          95 if language_missing else 100,
            "headings":          99 if headings_missing else 100,
            "tables":            AllyScores.score_from_count(tables_missing),
            "color_contrast":    AllyScores.contrast_score(contrast, total_text_runs),
            "lists":             AllyScores.score_from_count(lists),
            "links":             AllyScores.score_from_count(links),
        }

        final = min(scores.values())

        if final == 100:
            band = "Dark Green"
        elif final >= 67:
            band = "Light Green"
        elif final >= 34:
            band = "Yellow"
        else:
            band = "Red"

        scores["final"] = final
        scores["band"]  = band
        return scores
